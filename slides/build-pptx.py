#!/usr/bin/env python
# Build an editable .pptx from the reveal.js deck, faithful to its content & theme.
import sys
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SRC = "/Users/habib/Git/nodejs-2-days-training/slides/nodejs-training.html"
OUT = "/Users/habib/Git/nodejs-2-days-training/slides/nodejs-training.pptx"

INK   = RGBColor(0x17,0x21,0x2b)
MUTED = RGBColor(0x56,0x63,0x71)
TEAL  = RGBColor(0x0f,0x76,0x6e)
TEALL = RGBColor(0x99,0xf6,0xe4)
LINE  = RGBColor(0xd9,0xdd,0xd5)
PANEL = RGBColor(0xff,0xff,0xff)
PAPER = RGBColor(0xf7,0xf7,0xf4)
CODE  = RGBColor(0x11,0x18,0x27)
CODEF = RGBColor(0xe6,0xe8,0xee)
WHITE = RGBColor(0xff,0xff,0xff)
BLUEG = RGBColor(0xbf,0xcc,0xd7)
CALLBG= RGBColor(0xff,0xff,0xff)

EMU = 914400
SW, SH = 13.333, 7.5
ML, MR = 0.62, 0.62
CW = SW - ML - MR

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
blank = prs.slide_layouts[6]

def txt(el):
    for br in el.find_all("br"):
        br.replace_with("\n")
    return el.get_text().strip()

def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def textbox(slide, l, t, w, h, lines, size, color, bold=False, align=PP_ALIGN.LEFT,
            font="Inter", anchor=MSO_ANCHOR.TOP, line_spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        runs = ln if isinstance(ln, list) else [(ln, bold, color)]
        for j,(rt,rb,rc) in enumerate(runs):
            r = p.add_run(); r.text = rt
            r.font.size = Pt(size); r.font.bold = rb
            r.font.color.rgb = rc; r.font.name = font
    return tb

def rrect(slide, l, t, w, h, fill, line=None, radius=0.08):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp

def add_footer(slide, fl, fr, dark=False):
    c = BLUEG if dark else MUTED
    if fl: textbox(slide, ML, SH-0.5, CW*0.6, 0.3, fl, 9, c, bold=True)
    if fr: textbox(slide, ML+CW*0.4, SH-0.5, CW*0.6, 0.3, fr, 9, c, bold=True, align=PP_ALIGN.RIGHT)

def add_eyebrow(slide, text, dark=False):
    textbox(slide, ML, 0.5, CW, 0.3, text.upper(), 11, TEALL if dark else TEAL, bold=True)

# ---- block renderers (return new y) ----
def r_pills(slide, y, rows, dark):
    for row in rows:
        x = ML
        for label in row:
            w = min(0.16*len(label)+0.4, 3.2)
            if x + w > ML+CW:
                x = ML; y += 0.5
            bg = RGBColor(0x22,0x2f,0x3a) if dark else WHITE
            ln = RGBColor(0x3a,0x4a,0x57) if dark else LINE
            rrect(slide, x, y, w, 0.4, bg, line=ln, radius=0.5)
            textbox(slide, x, y, w, 0.4, label, 11, WHITE if dark else INK, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += w + 0.12
        y += 0.55
    return y + 0.1

def r_panels(slide, y, panels, cols, dark=False):
    gap = 0.22
    w = (CW - gap*(cols-1)) / cols
    h = 1.7 if any(p.get('big') for p in panels) else 1.85
    rows = [panels[i:i+cols] for i in range(0,len(panels),cols)]
    for row in rows:
        x = ML
        for p in row:
            rrect(slide, x, y, w, h, PANEL, line=LINE)
            iy = y + 0.16
            if p.get('big'):
                textbox(slide, x+0.16, iy, w-0.32, 0.7, p['big'], 34, TEAL, bold=True)
                iy += 0.75
            if p.get('h3'):
                textbox(slide, x+0.16, iy, w-0.32, 0.4, p['h3'], 14, INK, bold=True)
                iy += 0.42
            body = "\n".join(p.get('ps', []))
            if body:
                textbox(slide, x+0.16, iy, w-0.32, h-(iy-y)-0.14, body, 10.5, MUTED, line_spacing=1.12)
            x += w + gap
        y += h + 0.2
    return y

def r_table(slide, y, headers, rows):
    nr, nc = len(rows)+1, len(headers)
    h = 0.34 + 0.30*len(rows)
    gt = slide.shapes.add_table(nr, nc, Inches(ML), Inches(y), Inches(CW), Inches(h)).table
    # disable banding via style is complex; set cell colors manually
    for j,htext in enumerate(headers):
        c = gt.cell(0,j); c.text = htext
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xed,0xf3,0xf2)
        for pp in c.text_frame.paragraphs:
            for r in pp.runs: r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=INK; r.font.name="Inter"
        c.margin_top=Pt(2); c.margin_bottom=Pt(2)
    for i,row in enumerate(rows,1):
        for j,val in enumerate(row):
            c = gt.cell(i,j); c.text = val
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            for pp in c.text_frame.paragraphs:
                for r in pp.runs: r.font.size=Pt(10); r.font.color.rgb=MUTED; r.font.name="Inter"
            c.margin_top=Pt(2); c.margin_bottom=Pt(2)
    return y + h + 0.2

def r_code(slide, y, code):
    lines = code.split("\n")
    h = 0.22 + 0.205*len(lines)
    rrect(slide, ML, y, CW, h, CODE, radius=0.04)
    textbox(slide, ML+0.18, y+0.12, CW-0.36, h-0.24, lines, 11, CODEF,
            font="Consolas", line_spacing=1.05)
    return y + h + 0.2

def r_callout(slide, y, text):
    h = max(0.6, 0.32 + 0.26*(len(text)//95 + 1))
    rrect(slide, ML, y, 0.09, h, TEAL, radius=0.2)
    rrect(slide, ML+0.09, y, CW-0.09, h, CALLBG, line=LINE, radius=0.03)
    textbox(slide, ML+0.32, y, CW-0.5, h, text, 12.5, MUTED, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
    return y + h + 0.2

def r_flow(slide, y, nodes):
    x = ML; h = 0.55
    for n in nodes:
        if n['kind']=='arrow':
            textbox(slide, x, y, 0.35, h, n['t'], 13, MUTED, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += 0.4
        else:
            w = min(0.13*len(n['t'])+0.5, 3.4)
            if x+w > ML+CW: break
            rrect(slide, x, y, w, h, RGBColor(0xe6,0xf4,0xf1) if n.get('accent') else WHITE, line=LINE, radius=0.12)
            textbox(slide, x, y, w, h, n['t'], 11, INK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += w + 0.12
    return y + h + 0.25

def r_para(slide, y, text, dark=False):
    h = 0.3 + 0.26*(len(text)//110 + 1)
    textbox(slide, ML, y, CW, h, text, 13, BLUEG if dark else MUTED, line_spacing=1.2)
    return y + h + 0.12

# ---- parse ----
soup = BeautifulSoup(open(SRC, encoding="utf-8").read(), "lxml")
sections = soup.select(".reveal > .slides > section")
print(f"Parsing {len(sections)} sections")

for sec in sections:
    classes = sec.get("class", [])
    dark = ("title-slide" in classes) or ("chapter" in classes)
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK if dark else PAPER)

    eb = sec.select_one(".eyebrow")
    if eb: add_eyebrow(slide, txt(eb), dark)
    foot = sec.select(".footer span")
    fl = txt(foot[0]) if len(foot)>0 else ""
    fr = txt(foot[1]) if len(foot)>1 else ""
    add_footer(slide, fl, fr, dark)

    stack = sec.select_one(".center-stack") or sec.select_one(".slide-frame")
    y = 2.5 if dark else 1.95

    # walk direct children of stack
    for el in stack.find_all(recursive=False) if stack else []:
        cls = el.get("class", [])
        name = el.name
        if name == "h1":
            textbox(slide, ML, y, CW, 1.2, txt(el), 40, WHITE if dark else INK, bold=True, line_spacing=1.0)
            y += 1.15
        elif name == "h2":
            textbox(slide, ML, y, CW, 0.9, txt(el), 28, WHITE if dark else INK, bold=True, line_spacing=1.02)
            y += 0.95
        elif name == "p" and "subtitle" in cls:
            t = txt(el); h = 0.4 + 0.3*(len(t)//80 + 1)
            textbox(slide, ML, y, CW, h, t, 16, BLUEG if dark else MUTED, line_spacing=1.2)
            y += h + 0.15
        elif name == "p":
            y = r_para(slide, y, txt(el), dark)
        elif name == "div" and "meta" in cls:
            row = [txt(p) for p in el.select(".pill")]
            y = r_pills(slide, y, [row], dark)
        elif name == "div" and "grid" in cls:
            cols = 2 if "two" in cls else 3 if "three" in cls else 4 if "four" in cls else 5 if "five" in cls else 3
            panels=[]
            for pan in el.select(".panel"):
                d={}
                big = pan.select_one("strong")
                h3 = pan.select_one("h3")
                if "metric" in pan.get("class",[]) and big: d['big']=txt(big)
                if h3: d['h3']=txt(h3)
                d['ps']=[txt(p) for p in pan.find_all("p", recursive=False)]
                if not d.get('ps'):
                    d['ps']=[txt(p) for p in pan.select("p")]
                panels.append(d)
            y = r_panels(slide, y, panels, cols, dark)
        elif name == "div" and "flow" in cls:
            nodes=[]
            for ch in el.find_all(recursive=False):
                c=ch.get("class",[])
                if "arrow" in c: nodes.append({'kind':'arrow','t':txt(ch)})
                elif "node" in c: nodes.append({'kind':'node','t':txt(ch),'accent':'accent' in c})
            y = r_flow(slide, y, nodes)
        elif name == "div" and "callout" in cls:
            y = r_callout(slide, y, txt(el))
        elif name == "pre":
            code = el.select_one("code")
            y = r_code(slide, y, code.get_text().strip("\n"))
        elif name == "table":
            headers=[txt(th) for th in el.select("thead th")]
            rows=[[txt(td) for td in tr.select("td")] for tr in el.select("tbody tr")]
            y = r_table(slide, y, headers, rows)
        elif name == "ul":
            items=[txt(li) for li in el.select("li")]
            textbox(slide, ML, y, CW, 0.3*len(items)+0.2, ["• "+it for it in items], 13, MUTED, line_spacing=1.15)
            y += 0.32*len(items)+0.15

prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
